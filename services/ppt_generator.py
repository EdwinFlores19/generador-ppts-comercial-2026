from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
import os
import logging
from datetime import datetime
from lxml import etree

from services.scope_items import (
    EDITION_LABELS, get_scope_items, normalize_edition
)

log = logging.getLogger("ppt_generator")

# Configurar parser XML seguro contra vulnerabilidades XXE (XML External Entity Injection)
secure_parser = etree.XMLParser(
    resolve_entities=False,  # Evita expansión de entidades
    no_network=True,         # Evita conexiones de red para DTDs externas
    load_dtd=False,          # Evita carga de DTDs externas
    dtd_validation=False     # Deshabilita validación DTD
)
etree.set_default_parser(secure_parser)

# Paleta de Colores Corporativos de SEIDOR (extraída del tema oficial de la plantilla)
COLOR_PRIMARY = RGBColor(0x07, 0x15, 0x3A)     # Azul Noche SEIDOR (#07153A)
COLOR_ROYAL = RGBColor(0x26, 0x3C, 0x7A)       # Azul Real SEIDOR (#263C7A)
COLOR_SECONDARY = RGBColor(0x66, 0xB6, 0xFF)   # Azul Claro SEIDOR (#66B6FF)
COLOR_BACKGROUND = RGBColor(0xF6, 0xF6, 0xF6)  # Gris Perla SEIDOR (#F6F6F6)
COLOR_WHITE = RGBColor(255, 255, 255)          # Blanco Puro (#FFFFFF)
COLOR_TEXT = RGBColor(0x24, 0x25, 0x28)        # Gris Carbón de Texto (#242528)
COLOR_GRAY = RGBColor(0x91, 0x91, 0x91)        # Gris Neutro SEIDOR (#919191)
COLOR_CARD_LINE = RGBColor(0xD9, 0xE2, 0xF2)   # Contorno suave de tarjetas

# Tipografía corporativa oficial de la plantilla SEIDOR
FONT_HEADING = 'Arial'
FONT_BODY = 'Arial'

# Layouts oficiales usados por la plantilla "Capacitación de Joule":
# se buscan por nombre exacto (con preferencia de master) para heredar
# los fondos de ondas azules y el logo SEIDOR reales de la marca.
LAYOUT_SPEC_COVER = [(4, 'Logo Presentación Corporativa')]
LAYOUT_SPEC_CONTENT_WAVE = [(5, '1_Diseño grafico 2'), (7, '1_Diseño grafico 2')]
LAYOUT_SPEC_CONTENT_CLEAN = [(7, 'Blanca'), (5, 'Blanca')]
LAYOUT_SPEC_CLOSING = [(0, 'Cierre 1'), (0, 'Cierre 2')]

# Área segura de contenido (el logo SEIDOR del layout ocupa la esquina superior izquierda)
MARGIN_X = Inches(0.59)
CONTENT_WIDTH = Inches(12.15)
HEADER_TITLE_Y = Inches(1.02)
HEADER_SUB_Y = Inches(1.58)
CONTENT_TOP = Inches(2.05)
CONTENT_BOTTOM = Inches(6.9)


def clear_presentation_slides(prs):
    """
    Elimina todas las diapositivas de la presentación para comenzar con un lienzo limpio,
    preservando los Slide Masters, fuentes, logos y estilos corporativos.

    Parámetros:
    - prs (Presentation): Instancia de la presentación cargada de python-pptx.
    """
    id_list = prs.slides._sldIdLst
    for i in range(len(id_list) - 1, -1, -1):
        slide_id = id_list[i]
        r_id = slide_id.rId
        prs.part.drop_rel(r_id)
        del id_list[i]
    log.info("Se han limpiado todas las diapositivas existentes de la plantilla.")


def find_layout_by_keywords(prs, keywords):
    """
    Busca de forma flexible y robusta un patrón de diseño (Slide Layout)
    dentro de todos los Slide Masters de la presentación según una lista de palabras clave.

    Parámetros:
    - prs (Presentation): Instancia de la presentación cargada.
    - keywords (list): Lista de cadenas a buscar en minúsculas en el nombre del layout.

    Retorna:
    - SlideLayout: El primer layout que coincide con alguna de las palabras clave,
      o un layout por defecto de la presentación si no hay coincidencia.
    """
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            name_lower = layout.name.lower()
            for kw in keywords:
                if kw.lower() in name_lower:
                    log.info("Layout encontrado: '%s' usando la palabra clave '%s'", layout.name, kw)
                    return layout
    log.warning("No se encontró layout con palabras clave. Usando layout por defecto.")
    return prs.slide_layouts[0]


def find_layout_by_spec(prs, specs, fallback_keywords):
    """
    Localiza un layout por nombre exacto, respetando el índice de master preferido.
    Los nombres provienen de los layouts que usan las diapositivas reales de la
    plantilla corporativa de SEIDOR, garantizando el mismo fondo y estilo de marca.

    Parámetros:
    - prs (Presentation): Presentación cargada.
    - specs (list[tuple]): Lista de (índice_master, nombre_layout) en orden de preferencia.
    - fallback_keywords (list): Palabras clave de respaldo si ningún nombre coincide.
    """
    masters = list(prs.slide_masters)
    for master_idx, layout_name in specs:
        if master_idx < len(masters):
            for layout in masters[master_idx].slide_layouts:
                if layout.name == layout_name:
                    log.info("Layout de marca encontrado: master %d, '%s'", master_idx, layout_name)
                    return layout
    for _, layout_name in specs:
        for master in masters:
            for layout in master.slide_layouts:
                if layout.name == layout_name:
                    log.info("Layout de marca encontrado por nombre global: '%s'", layout_name)
                    return layout
    return find_layout_by_keywords(prs, fallback_keywords)


def remove_slide_placeholders(slide):
    """
    Elimina los marcadores de posición heredados del patrón de diseño
    para evitar interferencias con el renderizado personalizado.

    Parámetros:
    - slide (Slide): Diapositiva a limpiar.
    """
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            try:
                slide.shapes.element.remove(shape.element)
            except Exception as e:
                log.error("No se pudo remover el placeholder '%s': %s", shape.name, e)


MESES_ES = ["Enero", "Febrero", "Marzo", "Abril", "Mayo", "Junio",
            "Julio", "Agosto", "Septiembre", "Octubre", "Noviembre", "Diciembre"]


def _fecha_actual_es():
    """Retorna la fecha actual en formato 'Mes Año' en español (independiente del locale)."""
    now = datetime.now()
    return f"{MESES_ES[now.month - 1]} {now.year}"


def _set_text(paragraph, text, font=FONT_BODY, size=12, bold=False, color=COLOR_TEXT, align=PP_ALIGN.LEFT):
    """Aplica texto y formato completo a un párrafo."""
    paragraph.text = text
    paragraph.font.name = font
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    if align is not None:
        paragraph.alignment = align
    return paragraph


def add_header(slide, title_text, subtitle_text):
    """
    Agrega un encabezado estandarizado corporativo de SEIDOR a una diapositiva de contenido.
    El título comienza debajo del logo SEIDOR del layout para no superponerse con la marca.

    Parámetros:
    - slide (Slide): Diapositiva de destino.
    - title_text (str): Título principal de la lámina.
    - subtitle_text (str): Subtítulo o contexto de la lámina.
    """
    title_box = slide.shapes.add_textbox(MARGIN_X, HEADER_TITLE_Y, CONTENT_WIDTH, Inches(0.55))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_top = tf_title.margin_right = tf_title.margin_bottom = 0
    title_size = 22 if len(title_text) <= 60 else 19
    _set_text(tf_title.paragraphs[0], title_text, FONT_HEADING, title_size, True, COLOR_PRIMARY)

    sub_box = slide.shapes.add_textbox(MARGIN_X, HEADER_SUB_Y, CONTENT_WIDTH, Inches(0.4))
    tf_sub = sub_box.text_frame
    tf_sub.word_wrap = True
    tf_sub.margin_left = tf_sub.margin_top = tf_sub.margin_right = tf_sub.margin_bottom = 0
    sub_size = 12 if len(subtitle_text) <= 95 else 11
    _set_text(tf_sub.paragraphs[0], subtitle_text, FONT_BODY, sub_size, False, COLOR_ROYAL)


def _style_card(card, fill_color=COLOR_WHITE, line_color=COLOR_CARD_LINE, line_width=1.0):
    """Aplica el estilo visual de tarjeta corporativa SEIDOR a una forma."""
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    if line_color is None:
        card.line.fill.background()
    else:
        card.line.color.rgb = line_color
        card.line.width = Pt(line_width)
    card.shadow.inherit = False
    tf = card.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0.28)
    return tf


def style_table_header_cell(cell, text):
    """
    Aplica el diseño de celda de cabecera de tabla de SEIDOR.

    Parámetros:
    - cell (Cell): Celda a aplicar formato.
    - text (str): Texto a escribir en la celda.
    """
    cell.fill.solid()
    cell.fill.fore_color.rgb = COLOR_PRIMARY
    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.05)
    cell.margin_bottom = Inches(0.05)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set_text(cell.text_frame.paragraphs[0], text, FONT_HEADING, 12, True, COLOR_WHITE, PP_ALIGN.CENTER)


def style_table_cell(cell, text, is_even=False, bold=False, align=PP_ALIGN.LEFT, font_size=10):
    """
    Aplica el diseño de celda estándar de SEIDOR.

    Parámetros:
    - cell (Cell): Celda a formatear.
    - text (str): Texto a ingresar.
    - is_even (bool): Si la fila es par para alternar color de fondo.
    - bold (bool): Si el texto va en negrita.
    - align (PP_ALIGN): Alineación del texto.
    - font_size (int): Tamaño de fuente en puntos.
    """
    cell.fill.solid()
    cell.fill.fore_color.rgb = COLOR_BACKGROUND if is_even else COLOR_WHITE
    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.05)
    cell.margin_bottom = Inches(0.05)
    _set_text(cell.text_frame.paragraphs[0], text, FONT_BODY, font_size, bold, COLOR_TEXT, align)


def _add_slide_cover(prs, layout_cover, company_name, sector, edition):
    """
    Añade la portada corporativa usando el layout oficial de SEIDOR
    (fondo azul noche con ondas y logo centrado). El texto se ubica
    debajo del logo para respetar la composición de la marca.
    """
    slide = prs.slides.add_slide(layout_cover)
    remove_slide_placeholders(slide)

    ed = EDITION_LABELS[edition]
    kicker_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(12.13), Inches(0.4))
    tf_k = kicker_box.text_frame
    tf_k.word_wrap = True
    tf_k.margin_left = tf_k.margin_top = tf_k.margin_right = tf_k.margin_bottom = 0
    _set_text(tf_k.paragraphs[0], f"PROPUESTA COMERCIAL  |  {ed['programa'].upper()}",
              FONT_HEADING, 13, True, COLOR_SECONDARY)

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(4.95), Inches(12.13), Inches(1.9))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    tf1.margin_left = tf1.margin_right = tf1.margin_top = tf1.margin_bottom = 0

    _set_text(tf1.paragraphs[0], f"Propuesta de Transformación Digital: {ed['nombre']}",
              FONT_HEADING, 24, True, COLOR_WHITE, PP_ALIGN.CENTER)

    p_client = tf1.add_paragraph()
    client_size = 16 if len(company_name) > 35 else (18 if len(company_name) > 22 else 20)
    _set_text(p_client, f"Preparado para: {company_name}", FONT_HEADING, client_size, True,
              COLOR_SECONDARY, PP_ALIGN.CENTER)
    p_client.space_before = Pt(14)

    footer_box = slide.shapes.add_textbox(Inches(0.6), Inches(6.85), Inches(12.13), Inches(0.45))
    tf_f = footer_box.text_frame
    tf_f.word_wrap = True
    tf_f.margin_left = tf_f.margin_top = tf_f.margin_right = tf_f.margin_bottom = 0
    fecha = _fecha_actual_es()
    _set_text(tf_f.paragraphs[0],
              f"SEIDOR Consulting SAC  ·  Perú  ·  Sector {sector}  ·  {fecha}",
              FONT_BODY, 11, False, COLOR_WHITE, PP_ALIGN.CENTER)


def _add_slide_customer(prs, layout_content, sector, description):
    """Añade diapositiva de entendimiento del cliente."""
    slide = prs.slides.add_slide(layout_content)
    remove_slide_placeholders(slide)
    add_header(slide, "Entendimiento del Cliente y Contexto de Mercado",
               f"Operaciones en el Perú y tendencias clave en el sector: {sector}")

    card1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN_X, CONTENT_TOP, Inches(5.9), Inches(4.55))
    tf_c1 = _style_card(card1)
    _set_text(tf_c1.paragraphs[0], "Resumen del Cliente", FONT_HEADING, 17, True, COLOR_PRIMARY)
    p_c1_body = tf_c1.add_paragraph()
    _set_text(p_c1_body,
              f"{description}\n\nCon operaciones en territorio peruano, la compañía requiere optimizar su "
              "estructura funcional y contable para dar soporte a su crecimiento en el mediano y largo plazo.",
              FONT_BODY, 11 if len(description) > 220 else 12, False, COLOR_TEXT)
    p_c1_body.space_before = Pt(12)

    card2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.84), CONTENT_TOP, Inches(5.9), Inches(4.55))
    tf_c2 = _style_card(card2, fill_color=COLOR_BACKGROUND)
    _set_text(tf_c2.paragraphs[0], "Tendencias del Sector (2026)", FONT_HEADING, 17, True, COLOR_PRIMARY)
    p_c2_body = tf_c2.add_paragraph()
    _set_text(p_c2_body,
              f"Para mantener la competitividad en el mercado peruano de {sector}, el sector exige:\n\n"
              "• Integración total en la nube para procesos logísticos y financieros (SaaS).\n\n"
              "• Reportabilidad inmediata adaptada a las normativas y libros electrónicos de la SUNAT.\n\n"
              "• Adopción rápida de copilotos de Inteligencia Artificial (SAP Joule) para agilizar "
              "consultas transaccionales diarias.",
              FONT_BODY, 12, False, COLOR_TEXT)
    p_c2_body.space_before = Pt(12)


DEFAULT_PAINS = {
    "logistics": "Falta de trazabilidad y visibilidad en tiempo real del stock en almacenes. "
                 "Tiempos de compra prolongados y procesos de cotización manuales.",
    "financial": "Cierres contables mensuales lentos e ineficientes. Complejidad en la conciliación "
                 "de múltiples bancos y monedas extranjeras.",
    "management": "Inexistencia de un control de disponibilidad presupuestaria en tiempo real. "
                  "Silos de información desarticulados entre las distintas áreas operativas."
}


def _add_slide_pains(prs, layout_content, pains=None):
    """
    Añade diapositiva de dolores operativos.
    Si se reciben dolores personalizados (extraídos por el chatbot), se usan en
    lugar de los genéricos para personalizar la propuesta al prospecto.
    """
    slide = prs.slides.add_slide(layout_content)
    remove_slide_placeholders(slide)
    add_header(slide, "Dolores Operativos Clave", "Retos que frenan la eficiencia de la organización")

    merged = dict(DEFAULT_PAINS)
    if isinstance(pains, dict):
        for key in merged:
            val = pains.get(key)
            if isinstance(val, str) and val.strip():
                merged[key] = val.strip()

    pain_titles = ["1. Logística y Abastecimiento", "2. Gestión Contable y Financiera", "3. Control de Gestión Interno"]
    pain_desc = [merged["logistics"], merged["financial"], merged["management"]]

    card_width = Inches(3.883)
    card_height = Inches(4.55)
    for i in range(3):
        left_pos = Inches(0.59 + i * (3.883 + 0.25))
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, CONTENT_TOP, card_width, card_height)
        tf = _style_card(card, fill_color=COLOR_WHITE,
                         line_color=COLOR_SECONDARY if i == 0 else COLOR_CARD_LINE,
                         line_width=1.5 if i == 0 else 1.0)
        _set_text(tf.paragraphs[0], pain_titles[i], FONT_HEADING, 15, True, COLOR_PRIMARY)
        p_d = tf.add_paragraph()
        desc = pain_desc[i]
        _set_text(p_d, desc, FONT_BODY, 11.5 if len(desc) <= 260 else 10.5, False, COLOR_TEXT)
        p_d.space_before = Pt(14)


def _add_slide_grow(prs, layout_content, edition):
    """Añade diapositiva de solución estratégica según la edición (GROW/RISE)."""
    ed = EDITION_LABELS[edition]
    slide = prs.slides.add_slide(layout_content)
    remove_slide_placeholders(slide)
    add_header(slide, f"La Solución Estratégica: {ed['nombre']}",
               f"{ed['programa']}: acelerando el crecimiento operativo con mejores prácticas preconfiguradas")

    box_left = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN_X, CONTENT_TOP, Inches(5.9), Inches(4.55))
    tf_bl = _style_card(box_left, fill_color=COLOR_PRIMARY, line_color=None)
    _set_text(tf_bl.paragraphs[0], f"¿Por qué {ed['programa']}?", FONT_HEADING, 18, True, COLOR_WHITE)
    p_bl_b = tf_bl.add_paragraph()
    _set_text(p_bl_b,
              f"{ed['descripcion']}\n\n"
              "✓ Mitiga los dolores de stock mediante mejores prácticas integradas.\n\n"
              "✓ Acelera la localización contable peruana de forma nativa en la nube.\n\n"
              "✓ Habilita una plataforma segura y disponible 24/7 sin costos de infraestructura física local.",
              FONT_BODY, 12, False, COLOR_WHITE)
    p_bl_b.space_before = Pt(14)

    box_right = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.84), CONTENT_TOP, Inches(5.9), Inches(4.55))
    tf_br = _style_card(box_right, fill_color=COLOR_BACKGROUND)
    _set_text(tf_br.paragraphs[0], "Beneficios Directos para la Operación", FONT_HEADING, 18, True, COLOR_PRIMARY)
    p_br_b = tf_br.add_paragraph()
    _set_text(p_br_b,
              "• Experiencia de Usuario Moderna (Fiori): Procesos limpios e intuitivos que reducen "
              "drásticamente la curva de aprendizaje de los colaboradores.\n\n"
              "• Inteligencia Artificial (SAP Joule): Respuestas automáticas y comandos directos que "
              "eliminan la navegación engorrosa en menús complejos.\n\n"
              "• Trazabilidad Total: Trazabilidad contable instantánea desde la Solicitud de Pedido "
              "(SolPed) hasta el pago de la factura al proveedor.",
              FONT_BODY, 12, False, COLOR_TEXT)
    p_br_b.space_before = Pt(14)


def _add_scope_table(slide, header, row1, row2):
    """Construye la tabla de alcance funcional de dos módulos con estilo SEIDOR."""
    table_shape = slide.shapes.add_table(3, 3, MARGIN_X, CONTENT_TOP, Inches(12.15), Inches(4.5))
    table = table_shape.table
    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(3.5)
    table.columns[2].width = Inches(6.65)

    style_table_header_cell(table.cell(0, 0), header[0])
    style_table_header_cell(table.cell(0, 1), header[1])
    style_table_header_cell(table.cell(0, 2), header[2])

    for idx, row in enumerate((row1, row2), start=1):
        is_even = idx % 2 == 0
        style_table_cell(table.cell(idx, 0), row[0], is_even=is_even, bold=True, font_size=11)
        style_table_cell(table.cell(idx, 1), row[1], is_even=is_even, font_size=10.5)
        style_table_cell(table.cell(idx, 2), row[2], is_even=is_even, font_size=10.5)


def _add_slide_fi_mm(prs, layout_content):
    """Añade diapositiva de alcance FI y MM."""
    slide = prs.slides.add_slide(layout_content)
    remove_slide_placeholders(slide)
    add_header(slide, "Alcance Funcional Parte I: Procesos Core (FI & MM)",
               "Detalle de los módulos de Finanzas y Gestión de Materiales incluidos en el alcance")
    _add_scope_table(
        slide,
        ("Módulo Core", "Procesos de Negocio", "Detalle del Alcance e Impacto"),
        ("Finanzas (FI)",
         "• Estructura Organizativa\n• Contabilidad General\n• Cuentas por Pagar (CXP)\n"
         "• Cuentas por Cobrar (CXC)\n• Bancos y Activos Fijos\n• Cierre de Mes",
         "Configuración de sociedades y ledgers. Procesamiento de cobros y pagos con retenciones y "
         "detracciones oficiales del Perú. Gestión centralizada de bancos propios y conciliaciones "
         "electrónicas. Control contable de activos fijos (altas, depreciación, bajas). Reportes "
         "oficiales y libros electrónicos de la SUNAT."),
        ("Materiales (MM)",
         "• Estructura de Almacenes\n• Datos Maestros de Compras\n• Compras de Stock\n"
         "• Compras de Servicios\n• Gestión de Stocks",
         "Gestión de compras nacionales e importaciones con recargos y costos indirectos de transporte. "
         "Ciclo completo de adquisición desde la SolPed hasta la hoja de entrada de servicios (HES) y "
         "factura. Control y valorización de inventarios con movimientos de traspaso entre almacenes y "
         "regularizaciones.")
    )


def _add_slide_module_co_ps(prs, layout_content, complexity):
    """Añade diapositiva de alcance CO y PS solo si la complejidad es Alta."""
    if complexity != "Alta":
        return
    slide = prs.slides.add_slide(layout_content)
    remove_slide_placeholders(slide)
    add_header(slide, "Alcance Funcional Parte II: Control y Proyectos (CO & PS)",
               "Detalle de los procesos de Controlling y Project System de alta complejidad")
    _add_scope_table(
        slide,
        ("Módulo", "Procesos de Negocio", "Detalle del Alcance e Impacto"),
        ("Controlling (CO)",
         "• Centros de Costo (CECO)\n• Centros de Beneficio (CEBI)\n• Órdenes Internas de Gastos\n"
         "• Ciclos de Distribución\n• Análisis de Margen",
         "Estructura jerárquica de costos. Control exhaustivo de costos indirectos mediante colectores "
         "temporales (órdenes internas) y liquidaciones mensuales automáticas. Ciclos de distribución y "
         "subreparto de gastos generales. Análisis detallado de rentabilidad por segmento de mercado "
         "(CO-PA contable)."),
        ("Proyectos (PS)",
         "• Estructura de Proyectos (PEP)\n• Presupuesto y Disponibilidad\n• Imputaciones reales a PEP\n"
         "• Liquidación de Proyectos",
         "Estructuraciones WBS/PEP para controlar proyectos de gastos (OPEX) y proyectos de inversión "
         "(CAPEX). Control estricto de disponibilidad presupuestal (prevención de excesos de gasto). "
         "Liquidación mensual y cierre definitivo del proyecto contra activos fijos en curso o cuentas "
         "de balance.")
    )


MODULE_FULL_NAMES = {
    "FI": "Finanzas (FI)", "CO": "Controlling (CO)", "MM": "Materiales (MM)",
    "SD": "Ventas (SD)", "PP": "Producción (PP)", "PS": "Proyectos (PS)"
}


def _add_slide_scope_items(prs, layout_content, active_modules, edition):
    """
    Añade la lámina de Scope Items de SAP Best Practices para los módulos
    activos, según el catálogo editable de services/scope_items.py.
    """
    ed = EDITION_LABELS[edition]
    items_by_module = get_scope_items(active_modules)
    slide = prs.slides.add_slide(layout_content)
    remove_slide_placeholders(slide)
    add_header(slide, "Alcance Detallado: Scope Items SAP Best Practices",
               f"Procesos preconfigurados activados en {ed['nombre']}")

    rows = len(items_by_module) + 1
    table_height = Inches(min(4.6, 0.5 + rows * 0.6))
    table_shape = slide.shapes.add_table(rows, 2, MARGIN_X, CONTENT_TOP, Inches(12.15), table_height)
    table = table_shape.table
    table.columns[0].width = Inches(2.3)
    table.columns[1].width = Inches(9.85)
    style_table_header_cell(table.cell(0, 0), "Módulo")
    style_table_header_cell(table.cell(0, 1), "Scope Items de SAP Best Practices incluidos")

    # Formato compacto en párrafo continuo para que hasta 6 módulos
    # con todos sus scope items quepan en la lámina sin desbordes.
    for r, (mod, items) in enumerate(items_by_module, start=1):
        is_even = r % 2 == 0
        style_table_cell(table.cell(r, 0), MODULE_FULL_NAMES.get(mod, mod),
                         is_even=is_even, bold=True, font_size=10.5)
        parts = [f"{sid} {name}" if sid else name for sid, name in items]
        style_table_cell(table.cell(r, 1), "  ·  ".join(parts), is_even=is_even, font_size=9)

    note_box = slide.shapes.add_textbox(MARGIN_X, Inches(7.08), Inches(12.15), Inches(0.35))
    tf_n = note_box.text_frame
    tf_n.word_wrap = True
    tf_n.margin_left = tf_n.margin_top = tf_n.margin_right = tf_n.margin_bottom = 0
    _set_text(tf_n.paragraphs[0],
              "Scope items según la hoja S0 del Estimador SEIDOR (SAP Best Practices); "
              "el alcance definitivo se valida en la fase Explore (Fit-to-Standard).",
              FONT_BODY, 9, False, COLOR_GRAY)


def _add_slide_efficiency(prs, layout_content):
    """Añade diapositiva de comparación de eficiencias."""
    slide = prs.slides.add_slide(layout_content)
    remove_slide_placeholders(slide)
    add_header(slide, "Casos de Uso Prácticos / Matriz de Éxito",
               "Operación Tradicional vs Eficiencia en el Sistema SAP S/4HANA (Joule)")

    table_shape = slide.shapes.add_table(4, 2, MARGIN_X, CONTENT_TOP, Inches(12.15), Inches(4.5))
    table = table_shape.table
    table.columns[0].width = Inches(6.075)
    table.columns[1].width = Inches(6.075)
    style_table_header_cell(table.cell(0, 0), "Operación Tradicional (Dolor / Hoy)")
    style_table_header_cell(table.cell(0, 1), "Operación con SAP S/4HANA (Eficiencia Automática)")
    rows = [
        ("Búsqueda manual de facturas y documentos de proveedores navegando por múltiples menús o "
         "listando clientes de manera subjetiva e incompleta.",
         "Joule responde a comandos directos en lenguaje natural: \"Muéstrame las facturas de la "
         "sociedad 5710\" o \"Listar facturas vencidas\" en segundos, agilizando el flujo del analista."),
        ("Retrasos e imprecisiones al crear solicitudes de compra y órdenes sin políticas claras ni "
         "visualización inmediata del historial de precios y contratos.",
         "Acceso inmediato al historial completo de SolPeds y flujos de aprobación automáticos desde "
         "Fiori, con soporte integrado de Joule para ubicar contratos con mayor valor."),
        ("Silos de información desarticulados. Coordinaciones manuales ineficientes (ej. SAP vs Cloud) "
         "que conllevan a fricciones operativas y reprocesos constantes.",
         "Integración nativa completa de Finanzas, Compras y Ventas bajo el principio de \"Un Solo "
         "Equipo\" de SEIDOR, con trazabilidad 100% digital y una sola fuente de verdad."),
    ]
    for idx, (left_txt, right_txt) in enumerate(rows, start=1):
        is_even = idx % 2 == 0
        style_table_cell(table.cell(idx, 0), left_txt, is_even=is_even, font_size=11)
        style_table_cell(table.cell(idx, 1), right_txt, is_even=is_even, font_size=11)


def _add_slide_change(prs, layout_content):
    """Añade diapositiva de gestión del cambio."""
    slide = prs.slides.add_slide(layout_content)
    remove_slide_placeholders(slide)
    add_header(slide, "Gestión del Cambio y Cultura Organizacional",
               "Metodología de acompañamiento y adopción para el equipo de colaboradores")

    box_adop = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN_X, CONTENT_TOP, Inches(5.9), Inches(4.55))
    tf_ba = _style_card(box_adop, fill_color=COLOR_BACKGROUND)
    _set_text(tf_ba.paragraphs[0], "Acompañamiento en la Adopción", FONT_HEADING, 17, True, COLOR_PRIMARY)
    p_ba_b = tf_ba.add_paragraph()
    _set_text(p_ba_b,
              "La adopción tecnológica es liderada bajo un enfoque práctico:\n\n"
              "• Capacitación Focalizada: Talleres prácticos sobre el uso del chat de Joule, enseñando "
              "reglas claras (ej: evitar ambigüedades como 'clientes morosos' y usar términos exactos "
              "como 'facturas vencidas').\n\n"
              "• Fiori Quick-Wins: Capacitación intensiva en el portal Fiori para agilizar "
              "transacciones diarias.\n\n"
              "• Gestión Preventiva: Identificación y mitigación proactiva de la resistencia al cambio "
              "en las fases iniciales.",
              FONT_BODY, 11, False, COLOR_TEXT)
    p_ba_b.space_before = Pt(12)

    box_feed = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.84), CONTENT_TOP, Inches(5.9), Inches(4.55))
    tf_bf = _style_card(box_feed, fill_color=COLOR_PRIMARY, line_color=None)
    _set_text(tf_bf.paragraphs[0], "Liderazgo SEIDOR y Feedback Oportuno", FONT_HEADING, 17, True, COLOR_WHITE)
    p_bf_b = tf_bf.add_paragraph()
    _set_text(p_bf_b,
              "El equipo consultor de SEIDOR implementa la metodología de feedback estructurada para "
              "coordinar con los líderes:\n\n"
              "• Describir Situación: Enfocado en hechos objetivos y medibles del proyecto SAP, "
              "evitando juicios.\n\n"
              "• Impacto en Negocio: Explicar cómo afecta al cronograma o equipo (ej. retraso de FI "
              "afecta integrales).\n\n"
              '• Escucha Activa: Dar espacio al colaborador para entender su perspectiva ("¿Cómo lo ves tú?").\n\n'
              "• Acuerdos Concretos: Pactar acciones claras y medibles con plazos de cara a futuras entregas.",
              FONT_BODY, 11, False, COLOR_WHITE)
    p_bf_b.space_before = Pt(12)


def _add_slide_economics(prs, layout_content, summary):
    """Añade diapositiva de propuesta económica con tabla bimoneda y gráfico circular."""
    slide = prs.slides.add_slide(layout_content)
    remove_slide_placeholders(slide)
    add_header(slide, "Propuesta Económica e Inversión Localizada",
               f"Desglose de inversión bimoneda aplicable en el Perú (Tipo de Cambio: {summary['tipo_cambio_pen']})")

    table_shape_f = slide.shapes.add_table(7, 3, MARGIN_X, CONTENT_TOP, Inches(6.8), Inches(4.5))
    table_f = table_shape_f.table
    table_f.columns[0].width = Inches(3.2)
    table_f.columns[1].width = Inches(1.8)
    table_f.columns[2].width = Inches(1.8)
    style_table_header_cell(table_f.cell(0, 0), "Concepto de Inversión")
    style_table_header_cell(table_f.cell(0, 1), "USD")
    style_table_header_cell(table_f.cell(0, 2), "PEN")

    porcentaje_ams_label = int(round(summary.get('porcentaje_ams', 0.15) * 100))
    igv_pct = int(round(summary.get('factor_igv', 0.18) * 100))
    filas = [
        ("Licencias Anuales Cloud (SaaS)", 'licensing_cost_str', False),
        (f"Servicios de Implementación ({summary['total_hours']:.0f}h)", 'consulting_cost_str', False),
        (f"Soporte Anual Post Go-Live (AMS - {porcentaje_ams_label}%)", 'support_cost_str', False),
        ("Inversión Inicial Neta Año 1", 'net_investment_str', True),
        (f"Impuesto IGV ({igv_pct}%)", 'igv_str', True),
        ("Inversión Total Facturable", 'total_facturable_str', True),
    ]
    for r, (concepto, key, bold) in enumerate(filas, start=1):
        is_even = r % 2 == 0
        style_table_cell(table_f.cell(r, 0), concepto, is_even=is_even, bold=bold)
        style_table_cell(table_f.cell(r, 1), summary['usd'][key], is_even=is_even, bold=bold, align=PP_ALIGN.RIGHT)
        style_table_cell(table_f.cell(r, 2), summary['pen'][key], is_even=is_even, bold=bold, align=PP_ALIGN.RIGHT)

    chart_data = CategoryChartData()
    chart_data.categories = ['SaaS Cloud', 'Implementación', 'Soporte AMS']
    chart_data.add_series('Inversión', (summary['usd']['licensing_cost'],
                                        summary['usd']['consulting_cost'],
                                        summary['usd']['support_cost']))
    chart_shape = slide.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(7.7), CONTENT_TOP, Inches(5.0), Inches(4.5), chart_data)
    chart = chart_shape.chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "Composición de la Inversión Año 1 (USD)"
    title_p = chart.chart_title.text_frame.paragraphs[0]
    title_p.font.size = Pt(13)
    title_p.font.bold = True
    title_p.font.name = FONT_HEADING
    title_p.font.color.rgb = COLOR_PRIMARY
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(11)
    chart.legend.font.name = FONT_BODY

    series_pie = chart.series[0]
    slice_colors = [COLOR_SECONDARY, COLOR_PRIMARY, COLOR_GRAY]
    for idx, color in enumerate(slice_colors):
        try:
            point = series_pie.points[idx]
            fill = point.format.fill
            fill.solid()
            fill.fore_color.rgb = color
        except Exception as chart_err:
            log.error("Error al colorear porción del gráfico: %s", chart_err)

    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.show_percentage = True
    data_labels.show_value = False
    data_labels.number_format = '0%'
    data_labels.number_format_is_linked = False
    data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    data_labels.font.size = Pt(11)
    data_labels.font.bold = True
    data_labels.font.name = FONT_BODY
    data_labels.font.color.rgb = COLOR_TEXT


def _add_slide_roi(prs, layout_content, summary, exp_wks, real_wks, deploy_wks):
    """Añade diapositiva de cronograma y ROI con gráfico de columnas."""
    slide = prs.slides.add_slide(layout_content)
    remove_slide_placeholders(slide)
    add_header(slide, "Cronograma y Retorno de Inversión (ROI)",
               "Cronograma del proyecto (SAP Activate) y análisis de recuperación financiera")

    phases = ["Prepare", "Explore", "Realize", "Deploy", "Run"]
    durations = ["2 Semanas", f"{exp_wks:.1f} Semanas", f"{real_wks:.1f} Semanas", f"{deploy_wks:.1f} Semanas", "AMS"]
    tasks = ["Alineación y sandbox.", "Talleres de diseño y BPD.", "Configuración y pruebas.",
             "Migración y Go-Live.", "Soporte continuo."]
    width_block = Inches(2.26)
    height_block = Inches(1.55)
    y_timeline = Inches(1.95)
    for i in range(5):
        left_pos = Inches(0.59 + i * (2.26 + 0.21))
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, y_timeline, width_block, height_block)
        if i == 2:
            tf = _style_card(card, fill_color=COLOR_PRIMARY, line_color=None)
            text_color, title_color = COLOR_WHITE, COLOR_SECONDARY
        else:
            tf = _style_card(card, fill_color=COLOR_BACKGROUND)
            text_color, title_color = COLOR_TEXT, COLOR_PRIMARY
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0.12)
        _set_text(tf.paragraphs[0], phases[i], FONT_HEADING, 13, True, title_color)
        p_dur = tf.add_paragraph()
        _set_text(p_dur, durations[i], FONT_HEADING, 10, True, text_color)
        p_t = tf.add_paragraph()
        _set_text(p_t, tasks[i], FONT_BODY, 8.5, False, text_color)

    chart_data_col = CategoryChartData()
    roi_proj = summary['roi_projection']
    chart_data_col.categories = [f"Año {item['year']}" for item in roi_proj]
    chart_data_col.add_series('TCO Acumulado', tuple(item['cum_tco'] for item in roi_proj))
    chart_data_col.add_series('Ahorros Acumulados', tuple(item['cum_savings'] for item in roi_proj))
    chart_shape_col = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, MARGIN_X, Inches(3.75), Inches(8.2), Inches(3.15), chart_data_col)
    chart_col = chart_shape_col.chart
    chart_col.has_legend = True
    chart_col.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart_col.legend.include_in_layout = False
    chart_col.legend.font.size = Pt(10)
    chart_col.legend.font.name = FONT_BODY
    try:
        cat_axis = chart_col.category_axis
        cat_axis.tick_labels.font.size = Pt(10)
        cat_axis.tick_labels.font.name = FONT_BODY
        val_axis = chart_col.value_axis
        val_axis.tick_labels.font.size = Pt(9)
        val_axis.tick_labels.font.name = FONT_BODY
        val_axis.has_major_gridlines = True
    except Exception as axis_err:
        log.debug("No se pudo estilizar ejes del gráfico: %s", axis_err)
    try:
        fill_tco = chart_col.series[0].format.fill
        fill_tco.solid()
        fill_tco.fore_color.rgb = COLOR_PRIMARY
        fill_savings = chart_col.series[1].format.fill
        fill_savings.solid()
        fill_savings.fore_color.rgb = COLOR_SECONDARY
    except Exception as col_err:
        log.error("Error al colorear series de columnas: %s", col_err)

    box_metrics = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.0), Inches(3.75), Inches(3.74), Inches(3.15))
    tf_m = _style_card(box_metrics, fill_color=COLOR_PRIMARY, line_color=None)
    tf_m.margin_left = tf_m.margin_top = tf_m.margin_right = tf_m.margin_bottom = Inches(0.2)
    _set_text(tf_m.paragraphs[0], "Retorno de Inversión (TCO vs Ahorro)", FONT_HEADING, 14, True, COLOR_WHITE)
    anos_roi = summary.get('anos_roi', 5)
    p_m_b = tf_m.add_paragraph()
    _set_text(p_m_b,
              f"• Ahorro Anual Proyectado:\n  {summary['usd']['savings_annual_str']} USD\n"
              f"  ({summary['pen']['savings_annual_str']} PEN)\n"
              f"• Periodo de Recupero:\n  {summary['payback_period']:.2f} Años\n"
              f"• ROI Acumulado a {anos_roi} años:\n  {summary['roi_five_years']:.1f}%",
              FONT_BODY, 11, False, COLOR_WHITE)
    p_m_b.space_before = Pt(10)


def _add_slide_closing(prs, layout_closing):
    """
    Añade la diapositiva de cierre usando el layout oficial 'Cierre' de SEIDOR,
    que ya incluye el fondo azul noche y los datos de contacto de SEIDOR Perú.
    """
    slide = prs.slides.add_slide(layout_closing)
    remove_slide_placeholders(slide)
    t_box = slide.shapes.add_textbox(Inches(0.6), Inches(2.4), Inches(12.13), Inches(2.4))
    tf = t_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    _set_text(tf.paragraphs[0], "¡Muchas Gracias!", FONT_HEADING, 44, True, COLOR_WHITE, PP_ALIGN.CENTER)
    p_sub = tf.add_paragraph()
    _set_text(p_sub,
              "GROW with SAP: El futuro de la gestión empresarial en la nube, impulsado por SEIDOR Perú.",
              FONT_BODY, 16, False, COLOR_SECONDARY, PP_ALIGN.CENTER)
    p_sub.space_before = Pt(16)


def generate_deck(company_name, sector, description, complexity, financial_data, output_path,
                  pains=None, edition="Public"):
    """
    Genera la propuesta comercial en PowerPoint usando los layouts oficiales de la
    plantilla corporativa de SEIDOR (mismos fondos de ondas azules, logo y cierre).

    Parámetros:
    - company_name (str): Nombre del prospecto.
    - sector (str): Sector industrial.
    - description (str): Descripción del negocio del prospecto.
    - complexity (str): 'Alta' o 'Media' (define si se incluye la lámina CO & PS).
    - financial_data (dict): Resultado de financial_engine.calculate_financials().
    - output_path (str): Ruta del archivo .pptx de salida.
    - pains (dict, opcional): Dolores personalizados {'logistics','financial','management'}
      extraídos por el chatbot para personalizar la lámina de dolores.
    - edition (str, opcional): 'Public' (GROW with SAP) o 'Private' (RISE with SAP).
    """
    edition = normalize_edition(edition)
    template_name = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                                 "Capacitación de Joule - El futuro de SAP.pptx")
    if not os.path.exists(template_name):
        raise FileNotFoundError(f"La plantilla base corporativa no se encuentra en la ruta: {template_name}")
    prs = Presentation(template_name)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    clear_presentation_slides(prs)

    layout_cover = find_layout_by_spec(prs, LAYOUT_SPEC_COVER, ['separador secciones', 'portada', 'title slide'])
    layout_wave = find_layout_by_spec(prs, LAYOUT_SPEC_CONTENT_WAVE, ['blanca', 'plain', 'blank'])
    layout_clean = find_layout_by_spec(prs, LAYOUT_SPEC_CONTENT_CLEAN, ['blanca', 'plain', 'blank'])
    layout_closing = find_layout_by_spec(prs, LAYOUT_SPEC_CLOSING, ['cierre 1', 'cierre 2', 'cierre'])

    summary = financial_data['summary']
    modules = financial_data['modules']
    exp_wks = max([m['explore_weeks'] for m in modules.values()]) if modules else 5.68
    real_wks = max([m['realize_weeks'] for m in modules.values()]) if modules else 10.68
    deploy_wks = max([m['deploy_weeks'] for m in modules.values()]) if modules else 4.0

    _add_slide_cover(prs, layout_cover, company_name, sector, edition)
    _add_slide_customer(prs, layout_wave, sector, description)
    _add_slide_pains(prs, layout_wave, pains=pains)
    _add_slide_grow(prs, layout_wave, edition)
    _add_slide_fi_mm(prs, layout_clean)
    _add_slide_module_co_ps(prs, layout_clean, complexity)
    _add_slide_scope_items(prs, layout_clean, list(modules.keys()), edition)
    _add_slide_efficiency(prs, layout_clean)
    _add_slide_change(prs, layout_wave)
    _add_slide_economics(prs, layout_clean, summary)
    _add_slide_roi(prs, layout_clean, summary, exp_wks, real_wks, deploy_wks)
    _add_slide_closing(prs, layout_closing)
    prs.save(output_path)
    log.info("Presentación corporativa guardada con éxito en: %s", output_path)

if __name__ == "__main__":
    from services import financial_engine
    test_mods = ['FI', 'CO', 'MM', 'SD', 'PP', 'PS']
    f_data = financial_engine.calculate_financials(test_mods)
    generate_deck("Alicorp S.A.A.", "Consumo Masivo / Lácteos",
                  "Alicorp es una de las empresas lácteas y de consumo masivo más grandes del Perú.",
                  "Alta", f_data, "test_seidor_proposal.pptx")
