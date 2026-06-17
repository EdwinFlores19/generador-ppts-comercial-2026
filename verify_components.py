import os
import sqlite3
import scraper
import financial_engine
import ppt_generator
from pptx.enum.chart import XL_CHART_TYPE

DB_NAME = "proposals.db"

def run_tests():
    """
    Ejecuta el ciclo de pruebas de integración para validar la inyección de gráficos nativos,
    conversión multimoneda (USD/PEN), cálculo del IGV (18%) y generación del Pitch Deck.
    """
    print("==================================================")
    print("INICIANDO PRUEBAS DE INTEGRACIÓN ENTERPRISE")
    print("==================================================")
    
    # Test 1: Verificar Base de Datos SQLite y Localización
    print("\n--- Test 1: Verificación de Base de Datos y Parámetros de Localización ---")
    conn = sqlite3.connect(DB_NAME)
    cursor = conn.cursor()
    
    # Validar parámetros de localización peruana
    cursor.execute("SELECT parametro, valor FROM configuracion_comercial")
    rows = dict(cursor.fetchall())
    print("Ajustes comerciales y de localización en BD:", rows)
    
    assert 'factor_igv' in rows, "El parámetro 'factor_igv' debe existir en SQLite."
    assert rows['factor_igv'] == 0.18, "El IGV por defecto debe ser 0.18."
    assert 'tipo_cambio_pen' in rows, "El parámetro 'tipo_cambio_pen' debe existir en SQLite."
    assert rows['tipo_cambio_pen'] == 3.78, "El tipo de cambio por defecto debe ser 3.78."
    
    conn.close()
    print("Éxito: SQLite cuenta con los campos de localización y sus valores por defecto.")
    
    # Test 2: Verificar Scraper e Inteligencia Comercial
    print("\n--- Test 2: Verificación de Scraper y Heurística ---")
    gloria_profile = scraper.get_company_profile("Gloria S.A.")
    print("Perfil de Gloria S.A.:", gloria_profile)
    assert gloria_profile['complexity'] == 'Alta', "Gloria S.A. debe ser complejidad Alta por preset."
    
    # Test 3: Verificar Motor Financiero y Estructura Bimoneda con IGV
    print("\n--- Test 3: Verificación de Cálculos Bimoneda con IGV ---")
    active_mods = [m.strip() for m in gloria_profile['active_modules'].split(',')]
    fin_results = financial_engine.calculate_financials(active_mods)
    summary = fin_results['summary']
    
    # Validar presencia de USD y PEN
    assert 'usd' in summary, "El resumen debe incluir desglose en USD."
    assert 'pen' in summary, "El resumen debe incluir desglose en PEN."
    
    # Validar cálculos de IGV
    usd_net = summary['usd']['net_investment']
    usd_igv = summary['usd']['igv']
    usd_total = summary['usd']['total_facturable']
    assert round(usd_net * 0.18, 2) == usd_igv, f"Cálculo de IGV en USD incorrecto: {usd_net} * 0.18 != {usd_igv}"
    assert round(usd_net + usd_igv, 2) == usd_total, "La inversión total facturable en USD no coincide."
    
    # Validar tipo de cambio
    pen_net = summary['pen']['net_investment']
    assert round(usd_net * 3.78, 2) == pen_net, f"Conversión a PEN incorrecta: {usd_net} * 3.78 != {pen_net}"
    
    print("Éxito: Cálculos bimoneda e impacto del IGV validados matemáticamente.")
    
    # Test 4: Verificar Generación de PPTX y Gráficos Nativos
    print("\n--- Test 4: Verificación de Generación de PPTX y Gráficos Nativos ---")
    output_path = "generated_decks/Test_Gloria_Enterprise.pptx"
    os.makedirs("generated_decks", exist_ok=True)
    
    if os.path.exists(output_path):
        os.remove(output_path)
        
    ppt_generator.generate_deck(
        company_name="Leche Gloria S.A.",
        sector=gloria_profile['sector'],
        description=gloria_profile['description'],
        complexity=gloria_profile['complexity'],
        financial_data=fin_results,
        output_path=output_path
    )
    
    assert os.path.exists(output_path), "El archivo PPTX no fue generado."
    print("Éxito: Archivo PPTX guardado en:", output_path)
    
    # Test 5: Validar Contenido de Diapositivas y Gráficos
    print("\n--- Test 5: Validación de Elementos y Diapositivas ---")
    from pptx import Presentation
    prs = Presentation(output_path)
    print("Cantidad total de láminas generadas:", len(prs.slides))
    assert len(prs.slides) == 11, f"Se esperaban 11 diapositivas, se obtuvieron {len(prs.slides)}."
    
    # Comprobar presencia de gráficos (Slide 9 y Slide 10)
    # Slide 9 es la lámina 9 (índice 8)
    slide9 = prs.slides[8]
    has_chart9 = False
    for shape in slide9.shapes:
        if shape.has_chart:
            has_chart9 = True
            chart = shape.chart
            print(f"Lámina 9: Detectado gráfico nativo tipo: {chart.chart_type}")
            assert chart.chart_type == XL_CHART_TYPE.PIE, "El gráfico de la lámina 9 debe ser circular (PIE)."
            
    assert has_chart9, "La lámina 9 debe contener un gráfico nativo circular."
    
    # Slide 10 es la lámina 10 (índice 9)
    slide10 = prs.slides[9]
    has_chart10 = False
    for shape in slide10.shapes:
        if shape.has_chart:
            has_chart10 = True
            chart = shape.chart
            print(f"Lámina 10: Detectado gráfico nativo tipo: {chart.chart_type}")
            assert chart.chart_type == XL_CHART_TYPE.COLUMN_CLUSTERED, "El gráfico de la lámina 10 debe ser de columnas."
            
    assert has_chart10, "La lámina 10 debe contener un gráfico nativo de columnas."
    
    print("\n==================================================")
    print("¡TODAS LAS PRUEBAS ENTERPRISE PASARON CON ÉXITO!")
    print("==================================================")

if __name__ == "__main__":
    run_tests()
