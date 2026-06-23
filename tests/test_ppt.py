import os
import pytest
from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE
from services import financial_engine, ppt_generator


class TestPPTGenerator:
    @pytest.fixture(autouse=True)
    def setup(self):
        self.test_output = "generated_decks/test_pytest.pptx"
        os.makedirs("generated_decks", exist_ok=True)
        if os.path.exists(self.test_output):
            os.remove(self.test_output)

    def teardown_method(self):
        if os.path.exists(self.test_output):
            try:
                os.remove(self.test_output)
            except OSError:
                pass

    def test_module_constants(self):
        assert ppt_generator.COLOR_PRIMARY is not None
        assert ppt_generator.COLOR_SECONDARY is not None
        assert ppt_generator.FONT_HEADING is not None
        assert ppt_generator.FONT_BODY is not None

    def test_generate_deck_function_exists(self):
        assert callable(ppt_generator.generate_deck)

    def test_generates_pptx_file(self):
        modules = ['FI', 'CO', 'MM', 'SD', 'PP', 'PS']
        fin_data = financial_engine.calculate_financials(modules)
        ppt_generator.generate_deck(
            company_name="Test S.A.",
            sector="Servicios Comerciales",
            description="Empresa de prueba pytest.",
            complexity="Alta",
            financial_data=fin_data,
            output_path=self.test_output
        )
        assert os.path.exists(self.test_output), "PPTX no fue generado"

    def test_pptx_has_correct_slide_count(self):
        modules = ['FI', 'CO', 'MM', 'SD']
        fin_data = financial_engine.calculate_financials(modules)
        ppt_generator.generate_deck(
            company_name="Test S.A.",
            sector="Servicios Comerciales",
            description="Empresa de prueba.",
            complexity="Media",
            financial_data=fin_data,
            output_path=self.test_output
        )
        prs = Presentation(self.test_output)
        assert len(prs.slides) == 10, f"Esperado 10 slides, obtenido {len(prs.slides)}"

    def test_pptx_has_pie_chart(self):
        modules = ['FI', 'CO', 'MM', 'SD', 'PP', 'PS']
        fin_data = financial_engine.calculate_financials(modules)
        ppt_generator.generate_deck(
            company_name="Test S.A.",
            sector="Servicios Comerciales",
            description="Empresa de prueba para charts.",
            complexity="Alta",
            financial_data=fin_data,
            output_path=self.test_output
        )
        prs = Presentation(self.test_output)
        slide9 = prs.slides[8]
        has_pie = any(
            shape.has_chart and shape.chart.chart_type == XL_CHART_TYPE.PIE
            for shape in slide9.shapes
        )
        assert has_pie, "Slide 9 debería tener gráfico circular"

    def test_pptx_has_column_chart(self):
        modules = ['FI', 'CO', 'MM', 'SD', 'PP', 'PS']
        fin_data = financial_engine.calculate_financials(modules)
        ppt_generator.generate_deck(
            company_name="Test S.A.",
            sector="Servicios Comerciales",
            description="Empresa de prueba para charts.",
            complexity="Alta",
            financial_data=fin_data,
            output_path=self.test_output
        )
        prs = Presentation(self.test_output)
        slide10 = prs.slides[9]
        has_column = any(
            shape.has_chart and shape.chart.chart_type == XL_CHART_TYPE.COLUMN_CLUSTERED
            for shape in slide10.shapes
        )
        assert has_column, "Slide 10 debería tener gráfico de columnas"
