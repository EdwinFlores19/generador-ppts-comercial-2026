import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
import os
from lxml import etree

# Configurar parser XML seguro contra vulnerabilidades XXE (XML External Entity Injection)
secure_parser = etree.XMLParser(
    resolve_entities=False,  # Evita expansión de entidades
    no_network=True,         # Evita conexiones de red para DTDs externas
    load_dtd=False,          # Evita carga de DTDs externas
    dtd_validation=False     # Deshabilita validación DTD
)
etree.set_default_parser(secure_parser)

# Paleta de Colores Corporativos de SEIDOR
COLOR_PRIMARY = RGBColor(11, 35, 71)       # Azul Profundo SEIDOR (#0B2347)
COLOR_SECONDARY = RGBColor(0, 168, 235)    # Cian Eléctrico SEIDOR (#00A8EB)
COLOR_BACKGROUND = RGBColor(244, 246, 249)  # Gris Claro Neutro (#F4F6F9)
COLOR_WHITE = RGBColor(255, 255, 255)      # Blanco Puro (#FFFFFF)
COLOR_TEXT = RGBColor(51, 51, 51)          # Gris Oscuro de Texto (#333333)
COLOR_GRAY = RGBColor(108, 117, 125)       # Gris Frío de Contorno (#6C757D)

FONT_HEADING = 'Segoe UI'
FONT_BODY = 'Arial'

def clear_presentation_slides(prs):
    """
    Elimina todas las diapositivas de la presentación para comenzar con un lienzo limpio,
    preservando los Slide Masters, fuentes, logos y estilos corporativos.
    
    Parámetros:
    - prs (Presentation): Instancia de la presentación cargada de python-pptx.
    """
    id_list = prs.slides._sldIdLst
    for i in range(len(id_list) - 1, -1, -1):
        slide_id = id_list[i]
        rId = slide_id.rId
        prs.part.drop_rel(rId)
        del id_list[i]
    print("Se han limpiado todas las diapositivas existentes de la plantilla.")

def find_layout_by_keywords(prs, keywords):
    """
    Busca de forma flexible y robusta un patrón de diseño (Slide Layout)
    dentro de todos los Slide Masters de la presentación según una lista de palabras clave.
    
    Parámetros:
    - prs (Presentation): Instancia de la presentación cargada.
    - keywords (list): Lista de cadenas a buscar en minúsculas en el nombre del layout.
    
    Retorna:
    - SlideLayout: El primer layout que coincide con alguna de las palabras clave,
      o un layout por defecto de la presentación si no hay coincidencia.
    """
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            name_lower = layout.name.lower()
            for kw in keywords:
                if kw.lower() in name_lower:
                    print(f"Layout encontrado: '{layout.name}' usando la palabra clave '{kw}'")
                    return layout
    print("Advertencia: No se encontró layout con palabras clave. Usando layout por defecto.")
    return prs.slide_layouts[0]

def remove_slide_placeholders(slide):
    """
    Elimina los marcadores de posición heredados del patrón de diseño
    para evitar interferencias con el renderizado personalizado.
    
    Parámetros:
    - slide (Slide): Diapositiva a limpiar.
    """
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            try:
                slide.shapes.element.remove(shape.element)
            except Exception as e:
                print(f"No se pudo remover el placeholder '{shape.name}': {e}")

def add_header(slide, title_text, subtitle_text):
    """
    Agrega un encabezado estandarizado corporativo de SEIDOR a una diapositiva de contenido.
    
    Parámetros:
    - slide (Slide): Diapositiva de destino.
    - title_text (str): Título principal de la lámina.
    - subtitle_text (str): Subtítulo o contexto de la lámina.
    """
    # Cuadro de texto para el título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_top = tf_title.margin_right = tf_title.margin_bottom = 0
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = FONT_HEADING
    p_title.font.size = Pt(24)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_PRIMARY
    
    # Cuadro de texto para el subtítulo
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12.333), Inches(0.5))
    tf_sub = sub_box.text_frame
    tf_sub.word_wrap = True
    tf_sub.margin_left = tf_sub.margin_top = tf_sub.margin_right = tf_sub.margin_bottom = 0
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = subtitle_text
    p_sub.font.name = FONT_BODY
    
    # Ajuste de fuente dinámico para evitar desbordamiento en subtítulos largos
    if len(subtitle_text) > 85:
        p_sub.font.size = Pt(11)
    else:
        p_sub.font.size = Pt(13)
        
    p_sub.font.color.rgb = COLOR_GRAY

def style_table_header_cell(cell, text):
    """
    Aplica el diseño de celda de cabecera de tabla de SEIDOR.
    
    Parámetros:
    - cell (Cell): Celda a aplicar formato.
    - text (str): Texto a escribir en la celda.
    """
    cell.fill.solid()
    cell.fill.fore_color.rgb = COLOR_PRIMARY
    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.05)
    cell.margin_bottom = Inches(0.05)
    
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT_HEADING
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

def style_table_cell(cell, text, is_even=False, bold=False, align=PP_ALIGN.LEFT, font_size=10):
    """
    Aplica el diseño de celda estándar de SEIDOR.
    
    Parámetros:
    - cell (Cell): Celda a formatear.
    - text (str): Texto a ingresar.
    - is_even (bool): Si la fila es par para alternar color de fondo.
    - bold (bool): Si el texto va en negrita.
    - align (PP_ALIGN): Alineación del texto.
    - font_size (int): Tamaño de fuente en puntos.
    """
    cell.fill.solid()
    cell.fill.fore_color.rgb = COLOR_BACKGROUND if is_even else COLOR_WHITE
    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.05)
    cell.margin_bottom = Inches(0.05)
    
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = FONT_BODY
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = COLOR_TEXT

def generate_deck(company_name, sector, description, complexity, financial_data, output_path):
    """
    Genera una presentación comercial widescreen de 8 a 10 láminas utilizando la plantilla
    corporativa de SEIDOR como base, inyectando gráficos dinámicos circulares y de columnas
    junto con la localización bimoneda e impuestos locales.
    
    Parámetros:
    - company_name (str): Nombre de la empresa peruana.
    - sector (str): Sector industrial.
    - description (str): Resumen de la empresa.
    - complexity (str): Complejidad ('Alta' o 'Media').
    - financial_data (dict): Resumen financiero con bimoneda (USD/PEN) e IGV.
    - output_path (str): Ruta de destino para guardar el PPTX.
    """
    template_name = "Capacitación de Joule - El futuro de SAP.pptx"
    
    if not os.path.exists(template_name):
        raise FileNotFoundError(f"La plantilla base corporativa no se encuentra en la ruta: {template_name}")
        
    prs = Presentation(template_name)
    
    # Asegurar el tamaño widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Limpiar diapositivas previas de la plantilla
    clear_presentation_slides(prs)
    
    # Identificar layouts corporativos
    layout_cover = find_layout_by_keywords(prs, ['title slide', 'portada', 'título', 'diapositiva de título'])
    layout_content = find_layout_by_keywords(prs, ['blanca', 'plain', 'blank'])
    layout_separator = find_layout_by_keywords(prs, ['separador secciones', 'separador'])
    layout_closing = find_layout_by_keywords(prs, ['cierre 1', 'cierre 2', 'cierre'])
    
    # Extraer datos financieros requeridos
    summary = financial_data['summary']
    modules = financial_data['modules']
    
    # Duración de fases
    exp_wks = max([m['explore_weeks'] for m in modules.values()]) if modules else 5.68
    real_wks = max([m['realize_weeks'] for m in modules.values()]) if modules else 10.68
    deploy_wks = max([m['deploy_weeks'] for m in modules.values()]) if modules else 4.0
    
    # --- SLIDE 1: Portada Corporativa SEIDOR ---
    slide1 = prs.slides.add_slide(layout_cover)
    remove_slide_placeholders(slide1)
    
    title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12.333), Inches(3.0))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    tf1.margin_left = tf1.margin_right = tf1.margin_top = tf1.margin_bottom = 0
    
    p_title = tf1.paragraphs[0]
    p_title.text = "Propuesta de Transformación Digital:\nSAP S/4HANA Public Cloud"
    p_title.font.name = FONT_HEADING
    p_title.font.size = Pt(36)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_WHITE
    
    p_client = tf1.add_paragraph()
    p_client.text = f"\nPreparado para: {company_name}"
    p_client.font.name = FONT_HEADING
    
    if len(company_name) > 35:
        p_client.font.size = Pt(16)
    elif len(company_name) > 22:
        p_client.font.size = Pt(18)
    else:
        p_client.font.size = Pt(22)
        
    p_client.font.bold = True
    p_client.font.color.rgb = COLOR_SECONDARY
    
    footer_box = slide1.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(12.333), Inches(0.8))
    tf_f = footer_box.text_frame
    p_f = tf_f.paragraphs[0]
    p_f.text = "SEIDOR Perú  |  Lead Solution Architecture & Preventa SAP"
    p_f.font.name = FONT_BODY
    p_f.font.size = Pt(13)
    p_f.font.color.rgb = COLOR_WHITE
    
    # --- SLIDE 2: Entendimiento del Cliente ---
    slide2 = prs.slides.add_slide(layout_content)
    remove_slide_placeholders(slide2)
    add_header(slide2, "Entendimiento del Cliente y Contexto de Mercado", f"Operaciones en el Perú y tendencias clave en el sector: {sector}")
    
    card1 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.8), Inches(5.9), Inches(4.8))
    card1.fill.solid()
    card1.fill.fore_color.rgb = COLOR_BACKGROUND
    card1.line.color.rgb = COLOR_GRAY
    card1.line.width = Pt(1)
    
    tf_c1 = card1.text_frame
    tf_c1.word_wrap = True
    tf_c1.margin_left = tf_c1.margin_top = tf_c1.margin_right = tf_c1.margin_bottom = Inches(0.3)
    p_c1_title = tf_c1.paragraphs[0]
    p_c1_title.text = "Resumen del Cliente"
    p_c1_title.font.name = FONT_HEADING
    p_c1_title.font.size = Pt(18)
    p_c1_title.font.bold = True
    p_c1_title.font.color.rgb = COLOR_PRIMARY
    
    p_c1_body = tf_c1.add_paragraph()
    p_c1_body.text = f"\n{description}\n\nCon operaciones en territorio peruano, la compañía requiere optimizar su estructura funcional y contable para dar soporte a su crecimiento en el mediano y largo plazo."
    p_c1_body.font.name = FONT_BODY
    
    if len(description) > 220:
        p_c1_body.font.size = Pt(11)
    else:
        p_c1_body.font.size = Pt(12)
    p_c1_body.font.color.rgb = COLOR_TEXT
    
    card2 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.933), Inches(1.8), Inches(5.9), Inches(4.8))
    card2.fill.solid()
    card2.fill.fore_color.rgb = COLOR_BACKGROUND
    card2.line.color.rgb = COLOR_GRAY
    card2.line.width = Pt(1)
    
    tf_c2 = card2.text_frame
    tf_c2.word_wrap = True
    tf_c2.margin_left = tf_c2.margin_top = tf_c2.margin_right = tf_c2.margin_bottom = Inches(0.3)
    p_c2_title = tf_c2.paragraphs[0]
    p_c2_title.text = "Tendencias del Sector (2026)"
    p_c2_title.font.name = FONT_HEADING
    p_c2_title.font.size = Pt(18)
    p_c2_title.font.bold = True
    p_c2_title.font.color.rgb = COLOR_PRIMARY
    
    p_c2_body = tf_c2.add_paragraph()
    p_c2_body.text = f"\nPara mantener la competitividad en el mercado peruano de {sector}, el sector exige:\n\n" \
                     f"• Integración total en la nube para procesos logísticos y financieros (SaaS).\n\n" \
                     f"• Reportabilidad inmediata adaptada a las normativas y libros electrónicos de la SUNAT.\n\n" \
                     f"• Adopción rápida de copilotos de Inteligencia Artificial (SAP Joule) para agilizar consultas transaccionales diarias."
    p_c2_body.font.name = FONT_BODY
    p_c2_body.font.size = Pt(12)
    p_c2_body.font.color.rgb = COLOR_TEXT
    
    # --- SLIDE 3: Dolores Operativos ---
    slide3 = prs.slides.add_slide(layout_content)
    remove_slide_placeholders(slide3)
    add_header(slide3, "Dolores Operativos Clave", "Retos que frenan la eficiencia de la organización")
    
    pain_titles = [
        "1. Logística y Abastecimiento",
        "2. Gestión Contable y Financiera",
        "3. Control de Gestión Interno"
    ]
    pain_desc = [
        "Falta de trazabilidad y visibilidad en tiempo real del stock en almacenes.\nTiempos de compra prolongados y procesos de cotización manuales.",
        "Cierres contables mensuales lentos e ineficientes.\nComplejidad en la conciliación de múltiples bancos y monedas extranjeras.",
        "Inexistencia de un control de disponibilidad presupuestaria en tiempo real.\nSilos de información desarticulados entre las distintas áreas operativas."
    ]
    
    card_width = Inches(3.777)
    card_height = Inches(4.6)
    
    for i in range(3):
        left_pos = Inches(0.5 + i * (3.777 + 0.5))
        card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(1.8), card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_WHITE
        card.line.color.rgb = COLOR_SECONDARY if i == 0 else COLOR_GRAY
        card.line.width = Pt(1.5 if i == 0 else 1)
        
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0.25)
        
        p_t = tf.paragraphs[0]
        p_t.text = pain_titles[i]
        p_t.font.name = FONT_HEADING
        p_t.font.size = Pt(16)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_PRIMARY
        
        p_d = tf.add_paragraph()
        p_d.text = f"\n\n{pain_desc[i]}"
        p_d.font.name = FONT_BODY
        p_d.font.size = Pt(11.5)
        p_d.font.color.rgb = COLOR_TEXT
        
    # --- SLIDE 4: Solución GROW con SAP ---
    slide4 = prs.slides.add_slide(layout_content)
    remove_slide_placeholders(slide4)
    add_header(slide4, "La Solución Estratégica: SAP S/4HANA Public Cloud", "Acelerando el crecimiento operativo con mejores prácticas preconfiguradas")
    
    box_left = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.8), Inches(5.9), Inches(4.8))
    box_left.fill.solid()
    box_left.fill.fore_color.rgb = COLOR_PRIMARY
    box_left.line.fill.background()
    
    tf_bl = box_left.text_frame
    tf_bl.word_wrap = True
    tf_bl.margin_left = tf_bl.margin_top = tf_bl.margin_right = tf_bl.margin_bottom = Inches(0.3)
    p_bl_t = tf_bl.paragraphs[0]
    p_bl_t.text = "¿Por qué GROW con SAP S/4HANA?"
    p_bl_t.font.name = FONT_HEADING
    p_bl_t.font.size = Pt(20)
    p_bl_t.font.bold = True
    p_bl_t.font.color.rgb = COLOR_WHITE
    
    p_bl_b = tf_bl.add_paragraph()
    p_bl_b.text = f"\n\nGROW with SAP está diseñado específicamente para habilitar de forma rápida y escalable las operaciones de empresas en crecimiento.\n\n" \
                  f"✓ Mitiga los dolores de stock mediante mejores prácticas integradas.\n" \
                  f"✓ Acelera la localización contable peruana de forma nativa en la nube.\n" \
                  f"✓ Habilita una plataforma segura y disponible 24/7 sin costos de infraestructura física local."
    p_bl_b.font.name = FONT_BODY
    p_bl_b.font.size = Pt(12)
    p_bl_b.font.color.rgb = COLOR_WHITE
    
    box_right = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.933), Inches(1.8), Inches(5.9), Inches(4.8))
    box_right.fill.solid()
    box_right.fill.fore_color.rgb = COLOR_BACKGROUND
    box_right.line.color.rgb = COLOR_GRAY
    box_right.line.width = Pt(1)
    
    tf_br = box_right.text_frame
    tf_br.word_wrap = True
    tf_br.margin_left = tf_br.margin_top = tf_br.margin_right = tf_br.margin_bottom = Inches(0.3)
    p_br_t = tf_br.paragraphs[0]
    p_br_t.text = "Beneficios Directos para la Operación"
    p_br_t.font.name = FONT_HEADING
    p_br_t.font.size = Pt(20)
    p_br_t.font.bold = True
    p_br_t.font.color.rgb = COLOR_PRIMARY
    
    p_br_b = tf_br.add_paragraph()
    p_br_b.text = f"\n\n• Experiencia de Usuario Moderna (Fiori): Procesos limpios e intuitivos que reducen drásticamente la curva de aprendizaje de los colaboradores.\n\n" \
                  f"• Inteligencia Artificial (SAP Joule): Respuestas automáticas y comandos directos que eliminan la navegación engorrosa en menús complejos.\n\n" \
                  f"• Trazabilidad Total: Trazabilidad contable instantánea desde la Solicitud de Pedido (SolPed) hasta el pago de la factura al proveedor."
    p_br_b.font.name = FONT_BODY
    p_br_b.font.size = Pt(12)
    p_br_b.font.color.rgb = COLOR_TEXT
    
    # --- SLIDE 5: Alcance FI y MM ---
    slide5 = prs.slides.add_slide(layout_content)
    remove_slide_placeholders(slide5)
    add_header(slide5, "Alcance Funcional Parte I: Procesos Core (FI & MM)", "Detalle de los módulos de Finanzas y Gestión de Materiales incluidos en el alcance")
    
    rows = 3
    cols = 3
    left_t = Inches(0.5)
    top_t = Inches(1.8)
    width_t = Inches(12.333)
    height_t = Inches(4.5)
    
    table_shape = slide5.shapes.add_table(rows, cols, left_t, top_t, width_t, height_t)
    table = table_shape.table
    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(3.5)
    table.columns[2].width = Inches(6.833)
    
    style_table_header_cell(table.cell(0, 0), "Módulo Core")
    style_table_header_cell(table.cell(0, 1), "Procesos de Negocio")
    style_table_header_cell(table.cell(0, 2), "Detalle del Alcance e Impacto")
    
    style_table_cell(table.cell(1, 0), "Finanzas (FI)", is_even=False, bold=True, font_size=11)
    style_table_cell(table.cell(1, 1), "• Estructura Organizativa\n• Contabilidad General\n• Cuentas por Pagar (CXP)\n• Cuentas por Cobrar (CXC)\n• Bancos y Activos Fijos\n• Cierre de Mes", is_even=False, font_size=10.5)
    style_table_cell(table.cell(1, 2), "Configuración de sociedades y ledgers. Procesamiento de cobros y pagos con retenciones y detracciones oficiales del Perú. Gestión centralizada de bancos propios y conciliaciones electrónicas. Control contable de activos fijos (altas, depreciación, bajas). Reportes oficiales y libros electrónicos de la SUNAT.", is_even=False, font_size=10.5)
    
    style_table_cell(table.cell(2, 0), "Materiales (MM)", is_even=True, bold=True, font_size=11)
    style_table_cell(table.cell(2, 1), "• Estructura de Almacenes\n• Datos Maestros de Compras\n• Compras de Stock\n• Compras de Servicios\n• Gestión de Stocks", is_even=True, font_size=10.5)
    style_table_cell(table.cell(2, 2), "Gestión de compras nacionales e importaciones con recargos y costos indirectos de transporte. Ciclo completo de adquisición desde la SolPed hasta la hoja de entrada de servicios (HES) y factura. Control y valorización de inventarios con movimientos de traspaso entre almacenes y regularizaciones.", is_even=True, font_size=10.5)
    
    # --- SLIDE 6: Alcance CO y PS ---
    if complexity == "Alta":
        slide6 = prs.slides.add_slide(layout_content)
        remove_slide_placeholders(slide6)
        add_header(slide6, "Alcance Funcional Parte II: Control y Proyectos (CO & PS)", "Detalle de los procesos de Controlling y Project System de alta complejidad")
        
        table_shape2 = slide6.shapes.add_table(rows, cols, left_t, top_t, width_t, height_t)
        table2 = table_shape2.table
        table2.columns[0].width = Inches(2.0)
        table2.columns[1].width = Inches(3.5)
        table2.columns[2].width = Inches(6.833)
        
        style_table_header_cell(table2.cell(0, 0), "Módulo")
        style_table_header_cell(table2.cell(0, 1), "Procesos de Negocio")
        style_table_header_cell(table2.cell(0, 2), "Detalle del Alcance e Impacto")
        
        style_table_cell(table2.cell(1, 0), "Controlling (CO)", is_even=False, bold=True, font_size=11)
        style_table_cell(table2.cell(1, 1), "• Centros de Costo (CECO)\n• Centros de Beneficio (CEBI)\n• Órdenes Internas de Gastos\n• Ciclos de Distribución\n• Análisis de Margen", is_even=False, font_size=10.5)
        style_table_cell(table2.cell(1, 2), "Estructura jerárquica de costos. Control exhaustivo de costos indirectos mediante colectores temporales (órdenes internas) y liquidaciones mensuales automáticas. Ciclos de distribución y subreparto de gastos generales. Análisis detallado de rentabilidad por segmento de mercado (CO-PA contable).", is_even=False, font_size=10.5)
        
        style_table_cell(table2.cell(2, 0), "Proyectos (PS)", is_even=True, bold=True, font_size=11)
        style_table_cell(table2.cell(2, 1), "• Estructura de Proyectos (PEP)\n• Presupuesto y Disponibilidad\n• Imputaciones reales a PEP\n• Liquidación de Proyectos", is_even=True, font_size=10.5)
        style_table_cell(table2.cell(2, 2), "Estructuraciones WBS/PEP para controlar proyectos de gastos (OPEX) y proyectos de inversión (CAPEX). Control estricto de disponibilidad presupuestal (prevención de excesos de gasto). Liquidación mensual y cierre definitivo del proyecto contra activos fijos en curso o cuentas de balance.", is_even=True, font_size=10.5)
        
    # --- SLIDE 7: Comparación de Eficiencias ---
    slide7 = prs.slides.add_slide(layout_content)
    remove_slide_placeholders(slide7)
    add_header(slide7, "Casos de Uso Prácticos / Matriz de Éxito", "Operación Tradicional vs Eficiencia en el Sistema SAP S/4HANA (Joule)")
    
    rows_comp = 4
    cols_comp = 2
    
    table_shape3 = slide7.shapes.add_table(rows_comp, cols_comp, left_t, top_t, width_t, height_t)
    table3 = table_shape3.table
    table3.columns[0].width = Inches(6.166)
    table3.columns[1].width = Inches(6.166)
    
    style_table_header_cell(table3.cell(0, 0), "Operación Tradicional (Dolor / Hoy)")
    style_table_header_cell(table3.cell(0, 1), "Operación con SAP S/4HANA (Eficiencia Automática)")
    
    style_table_cell(table3.cell(1, 0), "Búsqueda manual de facturas y documentos de proveedores navegando por múltiples menús o listando clientes de manera subjetiva e incompleta.", is_even=False, font_size=11)
    style_table_cell(table3.cell(1, 1), "Joule responde a comandos directos en lenguaje natural: \"Muéstrame las facturas de la sociedad 5710\" o \"Listar facturas vencidas\" en segundos, agilizando el flujo del analista.", is_even=False, font_size=11)
    
    style_table_cell(table3.cell(2, 0), "Retrasos e imprecisiones al crear solicitudes de compra y órdenes sin políticas claras ni visualización inmediata del historial de precios y contratos.", is_even=True, font_size=11)
    style_table_cell(table3.cell(2, 1), "Acceso inmediato al historial completo de SolPeds y flujos de aprobación automáticos desde Fiori, con soporte integrado de Joule para ubicar contratos con mayor valor.", is_even=True, font_size=11)
    
    style_table_cell(table3.cell(3, 0), "Silos de información desarticulados. Coordinaciones manuales ineficientes (ej. SAP vs Cloud) que conllevan a fricciones operativas y reprocesos constantes.", is_even=False, font_size=11)
    style_table_cell(table3.cell(3, 1), "Integración nativa completa de Finanzas, Compras y Ventas bajo el principio de \"Un Solo Equipo\" de SEIDOR, con trazabilidad 100% digital y una sola fuente de verdad.", is_even=False, font_size=11)
    
    # --- SLIDE 8: Gestión del Cambio ---
    slide8 = prs.slides.add_slide(layout_content)
    remove_slide_placeholders(slide8)
    add_header(slide8, "Gestión del Cambio y Cultura Organizacional", "Metodología de acompañamiento y adopción para el equipo de colaboradores")
    
    box_adop = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.8), Inches(5.9), Inches(4.8))
    box_adop.fill.solid()
    box_adop.fill.fore_color.rgb = COLOR_BACKGROUND
    box_adop.line.color.rgb = COLOR_GRAY
    box_adop.line.width = Pt(1)
    
    tf_ba = box_adop.text_frame
    tf_ba.word_wrap = True
    tf_ba.margin_left = tf_ba.margin_top = tf_ba.margin_right = tf_ba.margin_bottom = Inches(0.3)
    
    p_ba_t = tf_ba.paragraphs[0]
    p_ba_t.text = "Acompañamiento en la Adopción"
    p_ba_t.font.name = FONT_HEADING
    p_ba_t.font.size = Pt(18)
    p_ba_t.font.bold = True
    p_ba_t.font.color.rgb = COLOR_PRIMARY
    
    p_ba_b = tf_ba.add_paragraph()
    p_ba_b.text = f"\n\nLa adopción tecnológica es liderada bajo un enfoque práctico:\n\n" \
                  f"• Capacitación Focalizada: Talleres prácticos sobre el uso del chat de Joule, enseñando reglas claras (ej: evitar ambigüedades como 'clientes morosos' y usar términos exactos como 'facturas vencidas').\n\n" \
                  f"• Fiori Quick-Wins: Capacitación intensiva en el portal Fiori para agilizar transacciones diarias.\n\n" \
                  f"• Gestión Preventiva: Identificación y mitigación proactiva de la resistencia al cambio en las fases iniciales."
    p_ba_b.font.name = FONT_BODY
    p_ba_b.font.size = Pt(11)
    p_ba_b.font.color.rgb = COLOR_TEXT
    
    box_feed = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.933), Inches(1.8), Inches(5.9), Inches(4.8))
    box_feed.fill.solid()
    box_feed.fill.fore_color.rgb = COLOR_PRIMARY
    box_feed.line.fill.background()
    
    tf_bf = box_feed.text_frame
    tf_bf.word_wrap = True
    tf_bf.margin_left = tf_bf.margin_top = tf_bf.margin_right = tf_bf.margin_bottom = Inches(0.3)
    
    p_bf_t = tf_bf.paragraphs[0]
    p_bf_t.text = "Liderazgo Seidor y Feedback Oportuno"
    p_bf_t.font.name = FONT_HEADING
    p_bf_t.font.size = Pt(18)
    p_bf_t.font.bold = True
    p_bf_t.font.color.rgb = COLOR_WHITE
    
    p_bf_b = tf_bf.add_paragraph()
    p_bf_b.text = f"\n\nEl equipo consultor de SEIDOR implementa la metodología de feedback estructurada para coordinar con los líderes:\n\n" \
                  f"• Describir Situación: Enfocado en hechos objetivos y medibles del proyecto SAP, evitando juicios.\n\n" \
                  f"• Impacto en Negocio: Explicar cómo afecta al cronograma o equipo (ej. retraso de FI afecta integrales).\n\n" \
                  f"• Escucha Activa: Dar espacio al colaborador para entender su perspectiva (\"¿Cómo lo ves tú?\").\n\n" \
                  f"• Acuerdos Concretos: Pactar acciones claras y medibles con plazos de cara a futuras entregas."
    p_bf_b.font.name = FONT_BODY
    p_bf_b.font.size = Pt(11)
    p_bf_b.font.color.rgb = COLOR_WHITE
    
    # --- SLIDE 9: Propuesta Económica e Inversión con IGV y Torta ---
    slide9 = prs.slides.add_slide(layout_content)
    remove_slide_placeholders(slide9)
    add_header(slide9, "Propuesta Económica e Inversión Localizada", f"Desglose de inversión bimoneda aplicable en el Perú (Tipo de Cambio: {summary['tipo_cambio_pen']})")
    
    # Tabla de desglose financiero bimoneda
    width_tf = Inches(6.8)
    height_tf = Inches(4.5)
    
    rows_f = 7
    cols_f = 3
    
    table_shape_f = slide9.shapes.add_table(rows_f, cols_f, Inches(0.5), Inches(1.8), width_tf, height_tf)
    table_f = table_shape_f.table
    table_f.columns[0].width = Inches(3.2)
    table_f.columns[1].width = Inches(1.8)
    table_f.columns[2].width = Inches(1.8)
    
    style_table_header_cell(table_f.cell(0, 0), "Concepto de Inversión")
    style_table_header_cell(table_f.cell(0, 1), "USD")
    style_table_header_cell(table_f.cell(0, 2), "PEN")
    
    # SaaS
    style_table_cell(table_f.cell(1, 0), "Licencias Anuales Cloud (SaaS)", is_even=False)
    style_table_cell(table_f.cell(1, 1), summary['usd']['licensing_cost_str'], is_even=False, align=PP_ALIGN.RIGHT)
    style_table_cell(table_f.cell(1, 2), summary['pen']['licensing_cost_str'], is_even=False, align=PP_ALIGN.RIGHT)
    
    # Consulting
    style_table_cell(table_f.cell(2, 0), f"Servicios de Implementación ({summary['total_hours']:.0f}h)", is_even=True)
    style_table_cell(table_f.cell(2, 1), summary['usd']['consulting_cost_str'], is_even=True, align=PP_ALIGN.RIGHT)
    style_table_cell(table_f.cell(2, 2), summary['pen']['consulting_cost_str'], is_even=True, align=PP_ALIGN.RIGHT)
    
    # Support (AMS)
    porcentaje_ams_label = int(summary.get('porcentaje_ams', 0.15) * 100) if 'porcentaje_ams' in summary else 15
    style_table_cell(table_f.cell(3, 0), f"Soporte Anual Post Go-Live (AMS - {porcentaje_ams_label}%)", is_even=False)
    style_table_cell(table_f.cell(3, 1), summary['usd']['support_cost_str'], is_even=False, align=PP_ALIGN.RIGHT)
    style_table_cell(table_f.cell(3, 2), summary['pen']['support_cost_str'], is_even=False, align=PP_ALIGN.RIGHT)
    
    # Inversión Neta
    style_table_cell(table_f.cell(4, 0), "Inversión Inicial Neta Año 1", is_even=True, bold=True)
    style_table_cell(table_f.cell(4, 1), summary['usd']['net_investment_str'], is_even=True, bold=True, align=PP_ALIGN.RIGHT)
    style_table_cell(table_f.cell(4, 2), summary['pen']['net_investment_str'], is_even=True, bold=True, align=PP_ALIGN.RIGHT)
    
    # IGV
    igv_pct = int(summary.get('factor_igv', 0.18) * 100)
    style_table_cell(table_f.cell(5, 0), f"Impuesto IGV ({igv_pct}%)", is_even=False, bold=True)
    style_table_cell(table_f.cell(5, 1), summary['usd']['igv_str'], is_even=False, bold=True, align=PP_ALIGN.RIGHT)
    style_table_cell(table_f.cell(5, 2), summary['pen']['igv_str'], is_even=False, bold=True, align=PP_ALIGN.RIGHT)
    
    # Total Facturable
    style_table_cell(table_f.cell(6, 0), "Inversión Total Facturable", is_even=True, bold=True)
    style_table_cell(table_f.cell(6, 1), summary['usd']['total_facturable_str'], is_even=True, bold=True, align=PP_ALIGN.RIGHT)
    style_table_cell(table_f.cell(6, 2), summary['pen']['total_facturable_str'], is_even=True, bold=True, align=PP_ALIGN.RIGHT)
    
    # Gráfico Circular (Pie Chart) Nativo a la derecha
    chart_data = CategoryChartData()
    chart_data.categories = ['SaaS Cloud', 'Implementación', 'Soporte AMS']
    chart_data.add_series('Inversión', (summary['usd']['licensing_cost'], summary['usd']['consulting_cost'], summary['usd']['support_cost']))
    
    x_pie = Inches(7.5)
    y_pie = Inches(1.8)
    cx_pie = Inches(5.3)
    cy_pie = Inches(4.5)
    
    chart_shape = slide9.shapes.add_chart(XL_CHART_TYPE.PIE, x_pie, y_pie, cx_pie, cy_pie, chart_data)
    chart = chart_shape.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    
    # Colores corporativos del gráfico circular
    series_pie = chart.series[0]
    slice_colors = [COLOR_SECONDARY, COLOR_PRIMARY, COLOR_GRAY]
    for idx, color in enumerate(slice_colors):
        try:
            point = series_pie.points[idx]
            fill = point.format.fill
            fill.solid()
            fill.fore_color.rgb = color
        except Exception as chart_err:
            print(f"Error al colorear porción del gráfico: {chart_err}")
            
    # --- SLIDE 10: Cronograma & ROI ---
    slide10 = prs.slides.add_slide(layout_content)
    remove_slide_placeholders(slide10)
    add_header(slide10, "Cronograma y Retorno de Inversión (ROI)", "Cronograma del proyecto (SAP Activate) y análisis de recuperación financiera")
    
    # Mover cronograma al tope y achicarlo
    phases = ["Prepare", "Explore", "Realize", "Deploy", "Run"]
    durations = ["2 Semanas", f"{exp_wks:.1f} Semanas", f"{real_wks:.1f} Semanas", f"{deploy_wks:.1f} Semanas", "AMS"]
    tasks = [
        "Alineación y sandbox.",
        "Talleres de diseño y BPD.",
        "Configuración y pruebas.",
        "Migración y Go-Live.",
        "Soporte continuo."
    ]
    
    width_block = Inches(2.26)
    height_block = Inches(1.7)
    y_timeline = Inches(1.7)
    
    for i in range(5):
        left_pos = Inches(0.5 + i * (2.26 + 0.25))
        card = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, y_timeline, width_block, height_block)
        card.fill.solid()
        if i == 2:
            card.fill.fore_color.rgb = COLOR_PRIMARY
            text_color = COLOR_WHITE
            title_color = COLOR_SECONDARY
        else:
            card.fill.fore_color.rgb = COLOR_BACKGROUND
            text_color = COLOR_TEXT
            title_color = COLOR_PRIMARY
            
        card.line.color.rgb = COLOR_GRAY
        card.line.width = Pt(1)
        
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0.12)
        
        p_ph = tf.paragraphs[0]
        p_ph.text = phases[i]
        p_ph.font.name = FONT_HEADING
        p_ph.font.size = Pt(13)
        p_ph.font.bold = True
        p_ph.font.color.rgb = title_color
        
        p_dur = tf.add_paragraph()
        p_dur.text = f"{durations[i]}\n"
        p_dur.font.name = FONT_HEADING
        p_dur.font.size = Pt(10)
        p_dur.font.bold = True
        p_dur.font.color.rgb = text_color
        
        p_t = tf.add_paragraph()
        p_t.text = tasks[i]
        p_t.font.name = FONT_BODY
        p_t.font.size = Pt(8.5)
        p_t.font.color.rgb = text_color
        
    # Gráfico de Columnas Nativo para la recuperación financiera (TCO vs Ahorros)
    chart_data_col = CategoryChartData()
    roi_proj = summary['roi_projection']
    chart_data_col.categories = [f"Año {item['year']}" for item in roi_proj]
    
    tco_values = [item['cum_tco'] for item in roi_proj]
    savings_values = [item['cum_savings'] for item in roi_proj]
    
    chart_data_col.add_series('TCO Acumulado', tuple(tco_values))
    chart_data_col.add_series('Ahorros Acumulados', tuple(savings_values))
    
    x_chart = Inches(0.5)
    y_chart = Inches(3.7)
    cx_chart = Inches(8.3)
    cy_chart = Inches(3.2)
    
    chart_shape_col = slide10.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x_chart, y_chart, cx_chart, cy_chart, chart_data_col)
    chart_col = chart_shape_col.chart
    chart_col.has_legend = True
    chart_col.legend.position = XL_LEGEND_POSITION.BOTTOM
    
    # Asignar colores corporativos al gráfico de columnas
    try:
        series_tco = chart_col.series[0]
        fill_tco = series_tco.format.fill
        fill_tco.solid()
        fill_tco.fore_color.rgb = COLOR_PRIMARY
        
        series_savings = chart_col.series[1]
        fill_savings = series_savings.format.fill
        fill_savings.solid()
        fill_savings.fore_color.rgb = COLOR_SECONDARY
    except Exception as col_err:
        print(f"Error al colorear series de columnas: {col_err}")
        
    # Tarjeta de Indicadores Financieros de Retorno
    box_metrics = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.3), Inches(3.7), Inches(3.5), Inches(3.2))
    box_metrics.fill.solid()
    box_metrics.fill.fore_color.rgb = COLOR_PRIMARY
    box_metrics.line.fill.background()
    
    tf_m = box_metrics.text_frame
    tf_m.word_wrap = True
    tf_m.margin_left = tf_m.margin_top = tf_m.margin_right = tf_m.margin_bottom = Inches(0.2)
    
    p_m_t = tf_m.paragraphs[0]
    p_m_t.text = "Retorno de Inversión (TCO vs Ahorro)"
    p_m_t.font.name = FONT_HEADING
    p_m_t.font.size = Pt(14)
    p_m_t.font.bold = True
    p_m_t.font.color.rgb = COLOR_WHITE
    
    anos_roi = summary.get('anos_roi', 5)
    p_m_b = tf_m.add_paragraph()
    p_m_b.text = f"\n• Ahorro Anual Proyectado:\n  {summary['usd']['savings_annual_str']} USD\n  ({summary['pen']['savings_annual_str']} PEN)\n" \
                  f"• Periodo de Recupero:\n  {summary['payback_period']:.2f} Años\n" \
                  f"• ROI Acumulado a {anos_roi} años:\n  {summary['roi_five_years']:.1f}%"
    p_m_b.font.name = FONT_BODY
    p_m_b.font.size = Pt(10)
    p_m_b.font.color.rgb = COLOR_WHITE
    
    # --- SLIDE 11: Cierre ---
    slide11 = prs.slides.add_slide(layout_closing)
    remove_slide_placeholders(slide11)
    
    t_box = slide11.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12.333), Inches(3.0))
    tf11 = t_box.text_frame
    tf11.word_wrap = True
    tf11.margin_left = tf11.margin_right = tf11.margin_top = tf11.margin_bottom = 0
    p_thank = tf11.paragraphs[0]
    p_thank.text = "¡Muchas Gracias!"
    p_thank.alignment = PP_ALIGN.CENTER
    p_thank.font.name = FONT_HEADING
    p_thank.font.size = Pt(44)
    p_thank.font.bold = True
    p_thank.font.color.rgb = COLOR_WHITE
    
    p_thank_sub = tf11.add_paragraph()
    p_thank_sub.text = "\nGROW with SAP: El futuro de la gestión empresarial en la nube, impulsado por SEIDOR Perú."
    p_thank_sub.alignment = PP_ALIGN.CENTER
    p_thank_sub.font.name = FONT_BODY
    p_thank_sub.font.size = Pt(16)
    p_thank_sub.font.color.rgb = COLOR_SECONDARY
    
    # Guardar presentación
    prs.save(output_path)
    print(f"Presentación corporativa guardada con éxito en: {output_path}")

if __name__ == "__main__":
    import financial_engine
    test_mods = ['FI', 'CO', 'MM', 'SD', 'PP', 'PS']
    f_data = financial_engine.calculate_financials(test_mods)
    generate_deck("Alicorp S.A.A.", "Consumo Masivo / Lácteos", "Alicorp es una de las empresas lácteas y de consumo masivo más grandes del Perú.", "Alta", f_data, "test_seidor_proposal.pptx")
